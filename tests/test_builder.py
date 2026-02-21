import os
import pytest
from pptx import Presentation
from src.builder import PPTXBuilder

@pytest.fixture
def mock_ir_data(tmp_path):
    """Generates a dummy JSON Intermediate Representation and a mock PNG."""
    # Create fake image
    dummy_img_path = os.path.join(tmp_path, "dummy_figure.png")
    with open(dummy_img_path, "wb") as f:
        f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82')

    return {
        "metadata": {
            "total_slides": 1,
            "aspect_ratio": 16/9,
            "page_count": 1,
        },
        "slides": [
            {
                "slide_index": 0,
                "image_width": 1920,
                "image_height": 1080,
                "background_color": "#FF5555",
                "elements": [
                    {
                        "type": "title",
                        "text": "Hello World",
                        "bbox": [100, 100, 500, 100],
                        "estimated_size": 48,
                        "is_bold": True
                    },
                    {
                        "type": "figure",
                        "source_file": dummy_img_path,
                        "bbox": [100, 300, 400, 400]
                    }
                ]
            }
        ]
    }

def test_slide_dimensions():
    """Slide canvas should always be the standard widescreen 13.333" x 7.5" regardless of DPI."""
    from src.builder import SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU
    builder = PPTXBuilder({}, "/tmp/dummy.pptx", dpi=300)
    assert builder.prs.slide_width  == SLIDE_WIDTH_EMU
    assert builder.prs.slide_height == SLIDE_HEIGHT_EMU

def test_hex_to_rgb():
    builder = PPTXBuilder({}, "/tmp/dummy.pptx")
    rgb = builder.hex_to_rgb("#FF0000")
    assert rgb == (255, 0, 0)
    
    rgb = builder.hex_to_rgb("00FF00")
    assert rgb == (0, 255, 0)

def test_builder_generation(mock_ir_data, tmp_path):
    out_path = os.path.join(tmp_path, "output.pptx")
    builder = PPTXBuilder(mock_ir_data, out_path, dpi=300)
    
    returned_path = builder.build()
    
    # Verify file was generated
    assert os.path.exists(returned_path)
    assert returned_path == out_path
    
    # Verify it can be opened by python-pptx
    prs = Presentation(returned_path)
    assert len(prs.slides) == 1
    
    slide = prs.slides[0]
    
    # Background should be set
    assert slide.background.fill.type is not None

    # Should have a shape for text and a shape for the picture
    assert len(slide.shapes) == 2
