import os
import logging
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from typing import Dict, Any

logger = logging.getLogger(__name__)

# Standard widescreen PowerPoint slide (16:9)
SLIDE_WIDTH_EMU  = 12192000   # 13.333 inches
SLIDE_HEIGHT_EMU = 6858000    # 7.5 inches


class PPTXBuilder:
    def __init__(self, ir_data: Dict[str, Any], output_path: str, dpi: int = 300):
        self.ir_data     = ir_data
        self.output_path = output_path
        self.dpi         = dpi
        self.prs         = Presentation()

        # Set standard widescreen canvas — fixed, independent of DPI
        self.prs.slide_width  = SLIDE_WIDTH_EMU
        self.prs.slide_height = SLIDE_HEIGHT_EMU

    def _make_scaler(self, img_w: int, img_h: int):
        """
        Returns (scale_x, scale_y) that map source-image pixels
        to EMUs on the standard slide canvas.

        Using the actual image pixel dimensions rather than DPI ensures
        correct scaling regardless of how the rasteriser was configured.
        """
        sx = SLIDE_WIDTH_EMU  / img_w
        sy = SLIDE_HEIGHT_EMU / img_h
        return sx, sy

    def _font_px_to_pt(self, px: float, scale_x: float) -> int:
        """
        Convert an estimated font size in source-image pixels to PowerPoint
        points, respecting the same spatial scale applied to coordinates.

        We use scale_x (the horizontal scale factor) to convert because font
        sizes are expressed in the same coordinate space as the image width.
        1 pt = 914400 / 72 EMU = 12700 EMU, so:
            pt = (px * scale_x) / 12700
        """
        emu = px * scale_x
        pt  = emu / 12700.0
        return max(6, int(pt))

    def hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Converts '#RRGGBB' to an RGBColor object."""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) != 6:
            return RGBColor(255, 255, 255)
        r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)

    def build(self) -> str:
        """
        Constructs the .pptx file from the IR dictionary and saves it.
        Returns the saved file path.
        """
        blank_layout = self.prs.slide_layouts[6]   # blank layout
        slides       = self.ir_data.get("slides", [])

        for slide_data in slides:
            logger.info(f"Building Slide {slide_data.get('slide_index', 0)}...")
            slide = self.prs.slides.add_slide(blank_layout)

            # Per-slide scale factors based on the source image size
            img_w = slide_data.get("image_width",  1920)
            img_h = slide_data.get("image_height", 1080)
            sx, sy = self._make_scaler(img_w, img_h)

            # Background colour
            bg_hex    = slide_data.get("background_color", "#FFFFFF")
            fill      = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self.hex_to_rgb(bg_hex)

            elements = slide_data.get("elements", [])
            for element in elements:
                el_type = element.get("type")
                bbox    = element.get("bbox")   # [x, y, w, h] in source pixels

                if not bbox or len(bbox) != 4:
                    continue

                x, y, w, h = bbox
                left   = int(x * sx)
                top    = int(y * sy)
                width  = max(1, int(w * sx))
                height = max(1, int(h * sy))

                if el_type == "figure":
                    source_file = element.get("source_file")
                    if source_file and os.path.exists(source_file):
                        try:
                            slide.shapes.add_picture(
                                source_file, left, top, width, height)
                        except Exception as e:
                            logger.error(f"Failed to add picture {source_file}: {e}")

                elif el_type in ("title", "text"):
                    text_str = element.get("text", "").strip()
                    if not text_str:
                        continue

                    textbox    = slide.shapes.add_textbox(left, top, width, height)
                    tf         = textbox.text_frame
                    tf.word_wrap = True

                    p      = tf.paragraphs[0]
                    p.text = text_str

                    est_px = element.get("estimated_size", 24)
                    pt     = self._font_px_to_pt(est_px, sx)

                    for run in p.runs:
                        run.font.size      = Pt(pt)
                        run.font.bold      = element.get("is_bold", False)
                        run.font.name      = "Arial"
                        run.font.color.rgb = RGBColor(0, 0, 0)

        logger.info(f"Saving final presentation to {self.output_path}")
        self.prs.save(self.output_path)
        logger.info(f"Presentation saved to: {self.output_path}")
        return self.output_path
